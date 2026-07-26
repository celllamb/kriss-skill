"""Rebuild a damaged PPTX into a clean python-pptx presentation.

This Level 3 repair favors openability and readable content over exact visual
fidelity. Animations, transitions, comments, custom layouts, notes, OLE
objects, exact positioning, and some charts may not be preserved.

Key principle: Passing ZIP/XML validation does not guarantee Microsoft
PowerPoint compatibility. PowerPoint is strict about the presentation
relationship graph, especially slideMaster and slideLayout relationships. It
is also strict about per-slide shape ID uniqueness and DrawingML paragraph end
run properties. If minimal XML/package repair is not enough, rebuild a clean
presentation graph.
"""

from __future__ import annotations

import json
import tempfile
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

try:  # pragma: no cover
    from inspect_pptx import OFFICE_REL, RelationshipRecord, inspect_pptx
    from validate_repaired_pptx import validate_pptx
except ImportError:  # pragma: no cover
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from inspect_pptx import OFFICE_REL, RelationshipRecord, inspect_pptx
    from validate_repaired_pptx import validate_pptx

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


@dataclass
class RebuildResult:
    """Result from a Level 3 rebuild."""

    input_path: str
    output_path: str
    warnings: list[str] = field(default_factory=list)
    extracted_slide_count: int = 0
    validation: dict[str, Any] | None = None

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


def rebuild_presentation(
    input_path: str | Path,
    output_path: str | Path,
    *,
    no_libreoffice: bool = False,
    strict: bool = False,
) -> RebuildResult:
    """Create a new PPTX with a clean master/layout graph and salvaged content."""

    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except ImportError as exc:
        raise RuntimeError(
            "Full rebuild requires python-pptx. Install it with `pip install python-pptx` "
            "or run a minimal/structural repair instead."
        ) from exc

    source = Path(input_path)
    output = Path(output_path)
    report = inspect_pptx(source)
    result = RebuildResult(
        input_path=str(source),
        output_path=str(output),
        warnings=[
            "Full rebuild preserves readable content where possible, not exact visual fidelity.",
            "Animations, transitions, comments, custom layouts, notes, OLE objects, exact positioning, and some charts may not be preserved.",
        ],
    )

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    slide_parts = report.slide_parts or _fallback_slide_order(source)
    result.extracted_slide_count = len(slide_parts)

    with zipfile.ZipFile(source, "r") as zf, tempfile.TemporaryDirectory(prefix="pptx-rebuild-") as temp_dir:
        temp_root = Path(temp_dir)
        for index, slide_part in enumerate(slide_parts, start=1):
            slide = prs.slides.add_slide(blank_layout)
            texts = _extract_texts(zf, slide_part)
            tables = _extract_tables(zf, slide_part)
            images = _extract_images(zf, slide_part, report.relationships, temp_root)

            left = Inches(0.55)
            top = Inches(0.45)
            width = prs.slide_width - Inches(1.1)
            if texts:
                box = slide.shapes.add_textbox(left, top, width, Inches(1.5))
                frame = box.text_frame
                frame.clear()
                for para_index, text in enumerate(texts[:24]):
                    paragraph = frame.paragraphs[0] if para_index == 0 else frame.add_paragraph()
                    paragraph.text = text
                    for run in paragraph.runs:
                        run.font.size = Pt(18 if para_index == 0 else 12)
                top = Inches(2.1)
            else:
                box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
                box.text_frame.text = f"Slide {index}"
                top = Inches(1.0)

            for table in tables[:2]:
                if not table:
                    continue
                rows = len(table)
                cols = max(len(row) for row in table)
                shape = slide.shapes.add_table(rows, cols, left, top, min(width, Inches(9)), Inches(0.35 * rows + 0.2))
                pptx_table = shape.table
                for r, row in enumerate(table):
                    for c, value in enumerate(row):
                        if c < cols:
                            pptx_table.cell(r, c).text = value
                top += Inches(0.4 * rows + 0.35)

            image_left = Inches(0.6)
            image_top = max(top, Inches(2.2))
            for image_path in images[:4]:
                try:
                    slide.shapes.add_picture(str(image_path), image_left, image_top, width=Inches(2.2))
                    image_left += Inches(2.35)
                    if image_left > Inches(8.0):
                        image_left = Inches(0.6)
                        image_top += Inches(1.8)
                except Exception as exc:  # pragma: no cover - depends on image codecs
                    result.warnings.append(f"Could not place image from {slide_part}: {exc}")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    validation = validate_pptx(
        output,
        expected_slide_count=result.extracted_slide_count,
        no_libreoffice=no_libreoffice,
        strict=strict,
    )
    result.validation = validation.to_dict()
    return result


def _fallback_slide_order(path: Path) -> list[str]:
    with zipfile.ZipFile(path, "r") as zf:
        return sorted(
            (name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")),
            key=lambda value: int("".join(ch for ch in Path(value).stem if ch.isdigit()) or "0"),
        )


def _extract_texts(zf: zipfile.ZipFile, slide_part: str) -> list[str]:
    try:
        root = ET.fromstring(zf.read(slide_part))
    except Exception:
        return []

    texts: list[str] = []
    for paragraph in root.findall(f".//{{{A_NS}}}p"):
        value = "".join(node.text or "" for node in paragraph.findall(f".//{{{A_NS}}}t")).strip()
        if value:
            texts.append(value)
    if not texts:
        value = "\n".join(node.text or "" for node in root.findall(f".//{{{A_NS}}}t")).strip()
        if value:
            texts.append(value)
    return texts


def _extract_tables(zf: zipfile.ZipFile, slide_part: str) -> list[list[list[str]]]:
    try:
        root = ET.fromstring(zf.read(slide_part))
    except Exception:
        return []

    tables: list[list[list[str]]] = []
    for tbl in root.findall(f".//{{{A_NS}}}tbl"):
        rows: list[list[str]] = []
        for tr in tbl.findall(f"{{{A_NS}}}tr"):
            row: list[str] = []
            for tc in tr.findall(f"{{{A_NS}}}tc"):
                row.append(" ".join((text.text or "") for text in tc.findall(f".//{{{A_NS}}}t")).strip())
            rows.append(row)
        if rows:
            tables.append(rows)
    return tables


def _extract_images(
    zf: zipfile.ZipFile,
    slide_part: str,
    relationships: list[RelationshipRecord],
    temp_root: Path,
) -> list[Path]:
    image_paths: list[Path] = []
    image_rels = [
        rel
        for rel in relationships
        if rel.source_part == slide_part
        and rel.relationship_type == f"{OFFICE_REL}/image"
        and rel.exists
        and rel.resolved_target
    ]
    for rel in image_rels:
        suffix = Path(rel.resolved_target or "").suffix or ".img"
        out_path = temp_root / f"{Path(slide_part).stem}_{rel.relationship_id}{suffix}"
        try:
            out_path.write_bytes(zf.read(rel.resolved_target or ""))
            image_paths.append(out_path)
        except Exception:
            continue
    return image_paths


def main(argv: list[str] | None = None) -> int:
    import argparse

    parser = argparse.ArgumentParser(description="Rebuild a damaged PPTX using python-pptx.")
    parser.add_argument("input", type=Path)
    parser.add_argument("output", type=Path)
    parser.add_argument("--strict", action="store_true")
    parser.add_argument("--no-libreoffice", action="store_true")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)

    try:
        result = rebuild_presentation(
            args.input,
            args.output,
            strict=args.strict,
            no_libreoffice=args.no_libreoffice,
        )
    except RuntimeError as exc:
        print(str(exc))
        return 2

    if args.json:
        print(json.dumps(result.to_dict(), indent=2))
    else:
        print(f"Wrote {result.output_path}")
        for warning in result.warnings:
            print(f"- {warning}")
    validation_ok = bool((result.validation or {}).get("ok"))
    return 0 if validation_ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
